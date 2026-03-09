#!/usr/bin/env python3
"""
lazy_pp.py — Generate PowerPoint from YAML.
Usage: python lazy_pp.py input.yaml [output.pptx]
"""

import sys
import os
import re
import yaml
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
from pptx.oxml.ns import qn
from pptx.chart.data import ChartData
from pptx.enum.chart import XL_CHART_TYPE
import copy
from lxml import etree


# ─────────────────────────── Helpers ────────────────────────────

def hex_to_rgb(h: str) -> RGBColor:
    h = h.lstrip("#")
    return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))

def inches(v):
    if v is None:
        return None
    return Inches(float(v))

def points(v):
    if v is None:
        return None
    return Pt(float(v))

ALIGN_MAP = {
    "left":    PP_ALIGN.LEFT,
    "center":  PP_ALIGN.CENTER,
    "right":   PP_ALIGN.RIGHT,
    "justify": PP_ALIGN.JUSTIFY,
}

CHART_MAP = {
    "bar":       XL_CHART_TYPE.COLUMN_CLUSTERED,
    "bar_stacked": XL_CHART_TYPE.COLUMN_STACKED,
    "hbar":      XL_CHART_TYPE.BAR_CLUSTERED,
    "line":      XL_CHART_TYPE.LINE,
    "pie":       XL_CHART_TYPE.PIE,
    "doughnut":  XL_CHART_TYPE.DOUGHNUT,
    "area":      XL_CHART_TYPE.AREA,
    "scatter":   XL_CHART_TYPE.XY_SCATTER,
}


# ─────────────────────────── Theme ─────────────────────────────

class Theme:
    def __init__(self, d: dict):
        self.bg         = d.get("background",   "FFFFFF")
        self.text       = d.get("text",          "1A1A1A")
        self.accent     = d.get("accent",        "0D5FAB")
        self.secondary  = d.get("secondary",     "6B7280")
        self.font_title = d.get("font_title",    "Calibri")
        self.font_body  = d.get("font_body",     "Calibri")
        self.title_size = d.get("title_size",    36)
        self.body_size  = d.get("body_size",     18)


# ────────────────────────── Slide builders ─────────────────────

def set_bg(slide, color: str):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = hex_to_rgb(color)


def add_textbox(slide, spec: dict, theme: Theme):
    """Generic text box element."""
    x  = inches(spec.get("x", 0.5))
    y  = inches(spec.get("y", 0.5))
    w  = inches(spec.get("w", 9))
    h  = inches(spec.get("h", 1))
    tf = slide.shapes.add_textbox(x, y, w, h).text_frame
    tf.word_wrap = spec.get("wrap", True)

    raw = spec.get("text", "")
    lines = raw if isinstance(raw, list) else [raw]

    for i, line in enumerate(lines):
        para = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        para.alignment = ALIGN_MAP.get(spec.get("align", "left"), PP_ALIGN.LEFT)

        run = para.add_run()
        if isinstance(line, dict):
            run.text = str(line.get("text", ""))
            opts = line
        else:
            run.text = str(line)
            opts = spec

        run.font.name  = opts.get("font",      theme.font_body)
        run.font.size  = points(opts.get("size", theme.body_size))
        run.font.bold  = opts.get("bold",   False)
        run.font.italic = opts.get("italic", False)
        color_key = opts.get("color", spec.get("color", theme.text))
        run.font.color.rgb = hex_to_rgb(color_key)

    space_after = spec.get("space_after")
    if space_after is not None:
        for para in tf.paragraphs:
            para.space_after = Pt(float(space_after))


def add_title_block(slide, title: str, subtitle: str | None, theme: Theme,
                    x=0.5, y=0.8, w=9.0, title_size=None, sub_size=None):
    ts = title_size or theme.title_size
    add_textbox(slide, {"x": x, "y": y, "w": w, "h": 1.2, "text": title,
                        "size": ts, "bold": True, "color": theme.text,
                        "font": theme.font_title, "align": "left"}, theme)
    if subtitle:
        add_textbox(slide, {"x": x, "y": y + 1.3, "w": w, "h": 0.7,
                            "text": subtitle, "size": sub_size or (theme.body_size - 2),
                            "color": theme.secondary, "font": theme.font_body,
                            "align": "left"}, theme)


def add_accent_bar(slide, theme: Theme, x=0.5, y=0.75, w=9.0, h=0.04):
    from pptx.util import Inches as I
    bar = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE = 1, but we need the enum val
        I(x), I(y), I(w), Pt(h * 72)
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = hex_to_rgb(theme.accent)
    bar.line.fill.background()


def add_bullets(slide, spec: dict, theme: Theme):
    x  = inches(spec.get("x", 0.5))
    y  = inches(spec.get("y", 0.5))
    w  = inches(spec.get("w", 9))
    h  = inches(spec.get("h", 4))
    tf = slide.shapes.add_textbox(x, y, w, h).text_frame
    tf.word_wrap = True

    items = spec.get("items", [])
    size  = spec.get("size", theme.body_size)

    for i, item in enumerate(items):
        if isinstance(item, dict):
            text  = str(item.get("text", ""))
            level = item.get("level", 0)
            bold  = item.get("bold", False)
            color = item.get("color", theme.text)
        else:
            text  = str(item)
            level = 0
            bold  = False
            color = theme.text

        para = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        para.level = level
        pPr = para._pPr
        if pPr is None:
            pPr = para._p.get_or_add_pPr()
        buChar = etree.SubElement(pPr, qn('a:buChar'))
        buChar.set('char', '•' if level == 0 else '–')

        indent_pt = 18 + level * 18
        pPr.set('indent', str(int(-indent_pt * 12700)))
        pPr.set('marL',   str(int( indent_pt * 12700)))

        run = para.add_run()
        run.text = text
        run.font.name   = spec.get("font", theme.font_body)
        run.font.size   = Pt(float(size))
        run.font.bold   = bold
        run.font.color.rgb = hex_to_rgb(color)

    space = spec.get("space_after", 4)
    for para in tf.paragraphs:
        para.space_after = Pt(float(space))


def add_image(slide, spec: dict):
    path = spec.get("path") or spec.get("file")
    if not path or not Path(path).exists():
        print(f"  [warn] image not found: {path}", file=sys.stderr)
        return
    x = inches(spec.get("x", 0.5))
    y = inches(spec.get("y", 0.5))
    w = inches(spec.get("w"))
    h = inches(spec.get("h"))
    kw = {}
    if w: kw["width"]  = w
    if h: kw["height"] = h
    slide.shapes.add_picture(path, x, y, **kw)


def add_shape(slide, spec: dict, theme: Theme):
    from pptx.enum.shapes import MSO_SHAPE_TYPE
    shape_map = {
        "rect":      1,   # MSO_CONNECTOR_TYPE not right; use add_shape
        "rectangle": 1,
        "oval":      9,
        "ellipse":   9,
        "rounded_rect": 5,
    }
    kind = spec.get("kind", "rectangle").lower()
    shape_type = shape_map.get(kind, 1)

    x  = inches(spec.get("x", 0.5))
    y  = inches(spec.get("y", 0.5))
    w  = inches(spec.get("w", 2))
    h  = inches(spec.get("h", 1))
    sh = slide.shapes.add_shape(shape_type, x, y, w, h)

    fill_color = spec.get("fill")
    if fill_color:
        sh.fill.solid()
        sh.fill.fore_color.rgb = hex_to_rgb(fill_color)
    else:
        sh.fill.background()

    line_color = spec.get("line_color")
    if line_color:
        sh.line.color.rgb = hex_to_rgb(line_color)
        sh.line.width     = Pt(spec.get("line_width", 1))
    else:
        sh.line.fill.background()

    label = spec.get("text")
    if label:
        tf = sh.text_frame
        tf.word_wrap = True
        run = tf.paragraphs[0].add_run()
        run.text = str(label)
        run.font.size  = points(spec.get("size", theme.body_size))
        run.font.color.rgb = hex_to_rgb(spec.get("color", theme.text))
        run.font.bold  = spec.get("bold", False)
        tf.paragraphs[0].alignment = ALIGN_MAP.get(spec.get("align", "center"), PP_ALIGN.CENTER)


def add_table(slide, spec: dict, theme: Theme):
    rows_data = spec.get("rows", [])
    if not rows_data:
        return
    nrows = len(rows_data)
    ncols = max(len(r) for r in rows_data)

    x = inches(spec.get("x", 0.5))
    y = inches(spec.get("y", 0.5))
    w = inches(spec.get("w", 9))
    h = inches(spec.get("h", nrows * 0.5))

    tbl = slide.shapes.add_table(nrows, ncols, x, y, w, h).table

    header_bg  = spec.get("header_bg",   theme.accent)
    header_fg  = spec.get("header_fg",   "FFFFFF")
    row_bg     = spec.get("row_bg",      "F8F8F8")
    alt_bg     = spec.get("alt_bg",      "FFFFFF")
    border_col = spec.get("border",      "D1D5DB")

    for ri, row in enumerate(rows_data):
        is_header = (ri == 0 and spec.get("header", True))
        for ci, cell_val in enumerate(row):
            if ci >= ncols:
                break
            cell = tbl.cell(ri, ci)
            cell.text = str(cell_val)

            tf   = cell.text_frame
            para = tf.paragraphs[0]
            run  = para.runs[0] if para.runs else para.add_run()
            run.text = str(cell_val)
            run.font.size  = points(spec.get("size", theme.body_size - 2))
            run.font.bold  = is_header
            run.font.name  = theme.font_body
            run.font.color.rgb = hex_to_rgb(header_fg if is_header else theme.text)
            para.alignment = PP_ALIGN.CENTER if is_header else PP_ALIGN.LEFT

            fill = cell.fill
            fill.solid()
            if is_header:
                fill.fore_color.rgb = hex_to_rgb(header_bg)
            elif ri % 2 == 0:
                fill.fore_color.rgb = hex_to_rgb(alt_bg)
            else:
                fill.fore_color.rgb = hex_to_rgb(row_bg)


def add_chart(slide, spec: dict, theme: Theme):
    # "chart_type" is the chart subtype; "type" is the element kind (="chart")
    chart_type_key = spec.get("chart_type", spec.get("kind", "bar")).lower()
    chart_type     = CHART_MAP.get(chart_type_key, XL_CHART_TYPE.COLUMN_CLUSTERED)

    cd = ChartData()
    categories = spec.get("categories", [])
    cd.categories = [str(c) for c in categories]

    series = spec.get("series", [])
    for s in series:
        cd.add_series(s.get("name", ""), s.get("values", []))

    x = inches(spec.get("x", 0.5))
    y = inches(spec.get("y", 0.5))
    w = inches(spec.get("w", 9))
    h = inches(spec.get("h", 4))

    slide.shapes.add_chart(chart_type, x, y, w, h, cd)


def add_line(slide, spec: dict, theme: Theme):
    """Horizontal or arbitrary line."""
    from pptx.util import Inches as I
    x1 = inches(spec.get("x",  0.5))
    y1 = inches(spec.get("y",  0.5))
    x2 = inches(spec.get("x2", spec.get("x", 0.5) + spec.get("w", 9)))
    y2 = inches(spec.get("y2", spec.get("y", 0.5)))

    connector = slide.shapes.add_connector(1, x1, y1, x2, y2)  # 1 = STRAIGHT
    connector.line.color.rgb = hex_to_rgb(spec.get("color", theme.secondary))
    connector.line.width     = Pt(spec.get("width", 1))


# ──────────────────────── Element dispatcher ────────────────────

def render_element(slide, el: dict, theme: Theme):
    kind = el.get("type", "text").lower()
    dispatch = {
        "text":    add_textbox,
        "textbox": add_textbox,
        "bullets": add_bullets,
        "image":   lambda s, e, t: add_image(s, e),
        "shape":   add_shape,
        "table":   add_table,
        "chart":   add_chart,
        "line":    add_line,
    }
    fn = dispatch.get(kind)
    if fn:
        fn(slide, el, theme)
    else:
        print(f"  [warn] unknown element type: {kind}", file=sys.stderr)


# ──────────────────────── Slide layouts ─────────────────────────

def build_slide_title(prs, spec: dict, theme: Theme):
    """Full-bleed title slide."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    set_bg(slide, spec.get("background", theme.bg))

    # Accent sidebar
    if spec.get("accent_bar", True):
        bar = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(0.25), Inches(5.625))
        bar.fill.solid()
        bar.fill.fore_color.rgb = hex_to_rgb(theme.accent)
        bar.line.fill.background()

    title    = spec.get("title",    "Untitled")
    subtitle = spec.get("subtitle", "")

    add_textbox(slide, {
        "x": 0.6, "y": 1.8, "w": 8.8, "h": 1.4,
        "text": title, "size": theme.title_size + 8,
        "bold": True, "color": spec.get("title_color", theme.text),
        "font": theme.font_title, "align": "left"
    }, theme)

    if subtitle:
        add_textbox(slide, {
            "x": 0.6, "y": 3.3, "w": 8.8, "h": 0.8,
            "text": subtitle, "size": theme.body_size,
            "color": theme.secondary, "font": theme.font_body, "align": "left"
        }, theme)

    for el in spec.get("elements", []):
        render_element(slide, el, theme)


def build_slide_content(prs, spec: dict, theme: Theme):
    """Standard content slide: title + body area."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, spec.get("background", theme.bg))

    title = spec.get("title", "")
    if title:
        add_textbox(slide, {
            "x": 0.5, "y": 0.3, "w": 9.0, "h": 0.7,
            "text": title, "size": theme.title_size,
            "bold": True, "color": theme.text,
            "font": theme.font_title, "align": "left"
        }, theme)
        # Thin accent line under title
        bar = slide.shapes.add_shape(1, Inches(0.5), Inches(1.0), Inches(9.0), Pt(2))
        bar.fill.solid()
        bar.fill.fore_color.rgb = hex_to_rgb(theme.accent)
        bar.line.fill.background()

    # Shorthand: bullets key at slide level
    if "bullets" in spec and "elements" not in spec:
        add_bullets(slide, {
            "x": 0.5, "y": 1.2, "w": 9.0, "h": 4.0,
            "items": spec["bullets"],
            "size": theme.body_size,
        }, theme)
    elif "content" in spec and "elements" not in spec:
        add_textbox(slide, {
            "x": 0.5, "y": 1.2, "w": 9.0, "h": 4.0,
            "text": spec["content"],
            "size": theme.body_size,
            "color": theme.text,
        }, theme)

    for el in spec.get("elements", []):
        render_element(slide, el, theme)


def build_slide_two_col(prs, spec: dict, theme: Theme):
    """Two-column layout."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, spec.get("background", theme.bg))

    title = spec.get("title", "")
    if title:
        add_textbox(slide, {
            "x": 0.5, "y": 0.3, "w": 9.0, "h": 0.7,
            "text": title, "size": theme.title_size,
            "bold": True, "color": theme.text, "font": theme.font_title
        }, theme)
        bar = slide.shapes.add_shape(1, Inches(0.5), Inches(1.0), Inches(9.0), Pt(2))
        bar.fill.solid()
        bar.fill.fore_color.rgb = hex_to_rgb(theme.accent)
        bar.line.fill.background()

    left  = spec.get("left",  {})
    right = spec.get("right", {})
    col_w = spec.get("col_width", 4.3)
    gap   = spec.get("gap",       0.4)
    y_start = spec.get("y_start", 1.2)

    for col_spec, x_off in [(left, 0.5), (right, 0.5 + col_w + gap)]:
        for el in col_spec.get("elements", []):
            el_copy = dict(el)
            el_copy["x"] = el_copy.get("x", 0) + x_off
            el_copy.setdefault("y", y_start)
            el_copy.setdefault("w", col_w)
            render_element(slide, el_copy, theme)

    for el in spec.get("elements", []):
        render_element(slide, el, theme)


def build_slide_section(prs, spec: dict, theme: Theme):
    """Section divider slide."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg = spec.get("background", theme.accent)
    set_bg(slide, bg)

    label = spec.get("label", "")
    if label:
        add_textbox(slide, {
            "x": 0.5, "y": 0.4, "w": 9, "h": 0.5,
            "text": label.upper(), "size": 12,
            "bold": True, "color": "FFFFFF",
            "font": theme.font_body, "align": "left"
        }, theme)

    add_textbox(slide, {
        "x": 0.5, "y": 1.5, "w": 9, "h": 2.0,
        "text": spec.get("title", ""),
        "size": theme.title_size + 6,
        "bold": True, "color": spec.get("text_color", "FFFFFF"),
        "font": theme.font_title
    }, theme)

    sub = spec.get("subtitle", "")
    if sub:
        add_textbox(slide, {
            "x": 0.5, "y": 3.6, "w": 9, "h": 0.8,
            "text": sub, "size": theme.body_size,
            "color": "DDDDDD", "font": theme.font_body
        }, theme)

    for el in spec.get("elements", []):
        render_element(slide, el, theme)


def build_slide_blank(prs, spec: dict, theme: Theme):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, spec.get("background", theme.bg))
    for el in spec.get("elements", []):
        render_element(slide, el, theme)


LAYOUT_BUILDERS = {
    "title":    build_slide_title,
    "content":  build_slide_content,
    "bullets":  build_slide_content,   # alias
    "two_col":  build_slide_two_col,
    "two-col":  build_slide_two_col,
    "section":  build_slide_section,
    "divider":  build_slide_section,
    "blank":    build_slide_blank,
}


# ─────────────────────────── Main ──────────────────────────────

def build(yaml_path: str, output_path: str):
    with open(yaml_path, "r") as f:
        doc = yaml.safe_load(f)

    # Presentation-level settings
    meta    = doc.get("meta", {})
    theme   = Theme(doc.get("theme", {}))
    slides  = doc.get("slides", [])

    prs = Presentation()
    prs.slide_width  = Inches(doc.get("width",  10))
    prs.slide_height = Inches(doc.get("height", 5.625))

    if meta.get("title"):   prs.core_properties.title   = meta["title"]
    if meta.get("author"):  prs.core_properties.author  = meta["author"]
    if meta.get("subject"): prs.core_properties.subject = meta["subject"]

    for i, slide_spec in enumerate(slides):
        layout = slide_spec.get("layout", "content").lower()
        builder = LAYOUT_BUILDERS.get(layout)
        if not builder:
            print(f"  [warn] slide {i+1}: unknown layout '{layout}', using 'content'",
                  file=sys.stderr)
            builder = build_slide_content
        builder(prs, slide_spec, theme)
        print(f"  slide {i+1:3d}: [{layout}] {slide_spec.get('title', '')}")

    prs.save(output_path)
    print(f"\n✓ saved → {output_path}")


def main():
    if len(sys.argv) < 2:
        print("Usage: python lazy_pp.py input.yaml [output.pptx]")
        sys.exit(1)

    yaml_path = sys.argv[1]
    if not Path(yaml_path).exists():
        print(f"Error: file not found: {yaml_path}", file=sys.stderr)
        sys.exit(1)

    stem        = Path(yaml_path).stem
    output_path = sys.argv[2] if len(sys.argv) > 2 else f"{stem}.pptx"

    print(f"Building: {yaml_path} → {output_path}")
    build(yaml_path, output_path)


if __name__ == "__main__":
    main()
